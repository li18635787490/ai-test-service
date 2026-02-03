"""
文档解析服务 - 支持多种文档格式
"""
import os
from pathlib import Path
from typing import Optional, Tuple
import aiofiles

from app.models import DocumentType


class DocumentParser:
    """文档解析器"""

    SUPPORTED_EXTENSIONS = {
        ".pdf": DocumentType.PDF,
        ".docx": DocumentType.DOCX,
        ".doc": DocumentType.DOCX,
        ".xlsx": DocumentType.XLSX,
        ".xls": DocumentType.XLSX,
        ".pptx": DocumentType.PPTX,
        ".ppt": DocumentType.PPTX,
        ".txt": DocumentType.TXT,
        ".md": DocumentType.MD,
    }

    @classmethod
    def get_document_type(cls, filename: str) -> Optional[DocumentType]:
        """根据文件名获取文档类型"""
        ext = Path(filename).suffix.lower()
        return cls.SUPPORTED_EXTENSIONS.get(ext)

    @classmethod
    async def parse(cls, file_path: str) -> Tuple[str, Optional[int]]:
        """
        解析文档内容

        Returns:
            Tuple[str, Optional[int]]: (文档内容, 页数)
        """
        doc_type = cls.get_document_type(file_path)

        if doc_type == DocumentType.PDF:
            return await cls._parse_pdf(file_path)
        elif doc_type == DocumentType.DOCX:
            return await cls._parse_docx(file_path)
        elif doc_type == DocumentType.XLSX:
            return await cls._parse_xlsx(file_path)
        elif doc_type == DocumentType.PPTX:
            return await cls._parse_pptx(file_path)
        elif doc_type in (DocumentType.TXT, DocumentType.MD):
            return await cls._parse_text(file_path)
        else:
            raise ValueError(f"不支持的文档格式: {file_path}")

    @classmethod
    async def _parse_pdf(cls, file_path: str) -> Tuple[str, int]:
        """解析 PDF 文档"""
        from PyPDF2 import PdfReader

        reader = PdfReader(file_path)
        page_count = len(reader.pages)

        content_parts = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text:
                content_parts.append(f"--- 第 {i + 1} 页 ---\n{text}")

        return "\n\n".join(content_parts), page_count

    @classmethod
    async def _parse_docx(cls, file_path: str) -> Tuple[str, None]:
        """解析 Word 文档"""
        from docx import Document

        doc = Document(file_path)

        content_parts = []

        # 提取段落
        for para in doc.paragraphs:
            if para.text.strip():
                # 检测标题样式
                if para.style.name.startswith('Heading'):
                    level = para.style.name[-1] if para.style.name[-1].isdigit() else '1'
                    content_parts.append(f"{'#' * int(level)} {para.text}")
                else:
                    content_parts.append(para.text)

        # 提取表格
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                table_text.append(row_text)
            if table_text:
                content_parts.append("\n[表格]\n" + "\n".join(table_text))

        return "\n\n".join(content_parts), None

    @classmethod
    async def _parse_xlsx(cls, file_path: str) -> Tuple[str, int]:
        """解析 Excel 文档"""
        from openpyxl import load_workbook

        wb = load_workbook(file_path, read_only=True, data_only=True)
        sheet_count = len(wb.sheetnames)

        content_parts = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            content_parts.append(f"--- 工作表: {sheet_name} ---")

            rows = []
            for row in sheet.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                if any(v.strip() for v in row_values):
                    rows.append(" | ".join(row_values))

            content_parts.append("\n".join(rows))

        wb.close()
        return "\n\n".join(content_parts), sheet_count

    @classmethod
    async def _parse_pptx(cls, file_path: str) -> Tuple[str, int]:
        """解析 PPT 文档"""
        from pptx import Presentation

        prs = Presentation(file_path)
        slide_count = len(prs.slides)

        content_parts = []

        for i, slide in enumerate(prs.slides):
            slide_content = [f"--- 幻灯片 {i + 1} ---"]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text)

            content_parts.append("\n".join(slide_content))

        return "\n\n".join(content_parts), slide_count

    @classmethod
    async def _parse_text(cls, file_path: str) -> Tuple[str, None]:
        """解析纯文本文档"""
        async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
            content = await f.read()
        return content, None
